#!/usr/bin/env python3
"""Build a .pptx deck summarizing the EUI icon search project.

The output is a self-contained .pptx file — uploading it to Google
Drive auto-converts it to Google Slides. Run this script first, then
upload via the Drive MCP (or `gdrive` CLI) and Drive does the rest.

Output: reports/eui_icon_search_deck.pptx
"""

from __future__ import annotations

import argparse
import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
PNG_DIR = REPO_ROOT / "reports" / "resvg_pngs_v115.0.0"

log = logging.getLogger("build_deck")

# Brand-ish palette (loosely Elastic-aligned without claiming brand fidelity)
ELASTIC_PINK = RGBColor(0xFE, 0xC5, 0x14)  # accent
ELASTIC_BLUE = RGBColor(0x00, 0x77, 0xCC)
DARK_NAVY = RGBColor(0x1A, 0x1C, 0x21)
MUTED = RGBColor(0x6F, 0x6F, 0x6F)
GREEN_OK = RGBColor(0x16, 0xA0, 0x65)
RED_BAD = RGBColor(0xC2, 0x39, 0x34)


# --- helpers ---------------------------------------------------------------


def add_title(slide, text: str) -> None:
    title = slide.shapes.title
    if title is None:
        return
    title.text = text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(32)
            r.font.bold = True
            r.font.color.rgb = DARK_NAVY


def add_bullets(slide, bullets: list[str | tuple[str, int]],
                left=Inches(0.6), top=Inches(1.6),
                width=Inches(11.5), height=Inches(5.5)) -> None:
    """Add a textbox of bullet points. Each bullet is either a string
    (level 0) or a (text, level) tuple."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.text = text
        for r in p.runs:
            r.font.size = Pt(20 if level == 0 else 16)
            r.font.color.rgb = DARK_NAVY if level == 0 else MUTED


def add_caption(slide, text: str, top=Inches(6.6),
                left=Inches(0.6), width=Inches(11.5)) -> None:
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    for r in p.runs:
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = MUTED


def add_image_with_label(slide, png: Path, label: str, left, top, side=Inches(2.0)) -> None:
    if png.exists():
        slide.shapes.add_picture(str(png), left, top, height=side, width=side)
    cap = slide.shapes.add_textbox(left, top + side + Inches(0.05), side, Inches(0.4))
    p = cap.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    for r in p.runs:
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED


def add_journey_bar(slide, stages: list[tuple[str, float]],
                    left_in: float = 0.7, top_in: float = 2.0,
                    width_in: float = 11.5, bar_h_in: float = 0.5) -> None:
    """Draw a horizontal bar for each stage showing top-1 percentage.

    All measurements as plain inches floats so the math stays simple;
    converted to EMU only when shapes are placed.
    """
    max_val = max(v for _, v in stages)
    label_w_in = 3.5
    bar_max_in = width_in - label_w_in - 1.5
    row_h_in = bar_h_in + 0.15
    for i, (label, pct) in enumerate(stages):
        y_in = top_in + row_h_in * i
        # label
        lbl = slide.shapes.add_textbox(
            Inches(left_in), Inches(y_in), Inches(label_w_in), Inches(bar_h_in)
        )
        p = lbl.text_frame.paragraphs[0]
        p.text = label
        for r in p.runs:
            r.font.size = Pt(14)
            r.font.color.rgb = DARK_NAVY
        # bar
        bw_in = bar_max_in * (pct / 100.0)
        is_max = pct == max_val
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_in + label_w_in), Inches(y_in),
            Inches(bw_in), Inches(bar_h_in),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ELASTIC_BLUE if is_max else RGBColor(0xB0, 0xD2, 0xEE)
        bar.line.fill.background()
        # value
        val = slide.shapes.add_textbox(
            Inches(left_in + label_w_in + bw_in + 0.1), Inches(y_in),
            Inches(1.3), Inches(bar_h_in),
        )
        vp = val.text_frame.paragraphs[0]
        vp.text = f"{pct:.1f}%"
        for r in vp.runs:
            r.font.size = Pt(14)
            r.font.bold = is_max
            r.font.color.rgb = DARK_NAVY


# --- slide builders -------------------------------------------------------


def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = "Multimodal Vector Search for EUI Icons"
    for r in p.runs:
        r.font.size = Pt(44)
        r.font.bold = True
        r.font.color.rgb = DARK_NAVY
    sub = s.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "From experiment to working tool — what we built, what we learned"
    for r in sp.runs:
        r.font.size = Pt(22)
        r.font.italic = True
        r.font.color.rgb = MUTED

    # 2. The problem
    s = prs.slides.add_slide(BLANK)
    add_title(s, "The Problem")
    add_bullets(s, [
        "EUI ships 677 icons. To use one you need its prop name (`<EuiIcon type=\"…\" />`).",
        "If you don't remember it, you scan the docs grid by eye.",
        "Goal: paste a screenshot, get the matching prop name back.",
        ("Existing demo used HuggingFace CLIP/MiniLM, ran offline, with no quality measurement.", 1),
    ])

    # 3. Stack
    s = prs.slides.add_slide(BLANK)
    add_title(s, "The Stack")
    add_bullets(s, [
        "Elasticsearch on Elastic Cloud — kNN over dense vectors",
        "jina-clip-v2 via Elastic Inference Service — Elastic-owned multimodal model, no API key, $0 incremental cost",
        "Express sidecar — holds the privileged ES key, normalizes via Sharp, returns ranked hits",
        "Docusaurus React component — `<IconSearch />` embedded in the EUI Icons docs page",
        "MCP server — AI assistants (Claude Code, Cursor) can call icon_search as a tool",
    ])

    # 4. How the search works
    s = prs.slides.add_slide(BLANK)
    add_title(s, "How a Search Runs")
    add_bullets(s, [
        "User pastes / drops an icon image (or types a description)",
        "Sidecar normalizes: flatten alpha onto white, resize-fit-contain to 256×256 (lanczos)",
        "Embed via `_inference/embedding/eui-icon-encoder` → 1024-dim vector",
        "ES kNN against `image_vector_aug_centroid`, version-filtered",
        "Dedupe by asset_filename (handles deprecation aliases like `merge` ↔ `submodule`)",
        "Return top-12 ranked candidates with cosine scores",
    ])

    # 5. Index coverage
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Index Coverage")
    add_bullets(s, [
        "6 EUI versions ingested: v91, v95, v100, v105, v110, v115",
        "Spans 24 minor releases of EUI history (~5 years)",
        "3,211 total docs — every doc has both a legacy `image_vector` and an augmented centroid",
        "New EUI versions ingestable in ~60 sec each via `python -m ingester run --version vX.Y.Z`",
    ])

    # 6. Quality discipline
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Measuring Quality Honestly")
    add_bullets(s, [
        "Quality has to be measured against what users actually paste — not the same image you indexed.",
        "Two test corpora:",
        ("Canonical (resvg) — every SVG rasterized server-side", 1),
        ("Browser-rendered (Playwright) — real screenshots from the live docs grid", 1),
        "Sweep: query with Playwright PNG → kNN → alias-aware rank lookup → report",
        "Caught two early bugs:",
        ("Leaky train/test split inflating accuracy from ~53% to a misleading 99%", 1),
        ("Missing Sharp normalize step in the test path (test ≠ production)", 1),
    ])

    # 7. The journey
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Quality Journey: +33pp top-1 in one day")
    add_journey_bar(s, [
        ("Baseline (raw bytes → ES)", 52.7),
        ("+ augmented centroids", 79.2),
        ("+ logo fillNegative fix", 78.9),
        ("+ token chrome fix", 84.7),
        ("+ app fillSecondary fix", 85.2),
        ("+ Borealis palette fix", 85.9),
    ])
    add_caption(s, "v115.0.0 held-out evaluation — Playwright queries, alias-aware ranking, score against 615 evaluable icons")

    # 8. Augmented centroids — the biggest single win
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Augmented Per-Icon Centroids — the biggest win (+26.5pp)")
    add_bullets(s, [
        "Naïve: index one PNG per icon, kNN against that one vector",
        "Problem: real user pastes have varying padding/cropping; tight match against a single render is brittle",
        "Insight: augment at INDEX time, not query time",
        ("Render each icon at 4 padding variants (0/10/25/50%)", 1),
        ("Embed each variant separately", 1),
        ("Mean-pool the 4 vectors as the icon's stored representation", 1),
        "No training. No extra storage at query time. Just smarter index construction.",
    ])

    # 9. Visual mismatch story (with icon comparisons if available)
    s = prs.slides.add_slide(BLANK)
    add_title(s, "The Visual Mismatch Detective Story")
    add_bullets(s, [
        "Pattern: every \"this looks wrong\" observation was a categorical pipeline bug affecting 12-60 icons.",
        "euiIcon__fillNegative (19 logos) — invisible dark outlines. logoCloud rank #19 → #1.",
        "Token chrome (57 tokens) — indexed as B&W glyphs, displayed as colored chips. 21 tokens jumped to #1.",
        "euiIcon__fillSecondary (60 App icons) — missing blue accent. codeApp / metricsApp miss → #1.",
        "Amsterdam vs Borealis palette — same color names, different hex values for v110+.",
    ])

    # Add small comparison strip if PNGs available
    if PNG_DIR.exists():
        try:
            for i, prop in enumerate(["logoCloud", "codeApp", "tokenParameter"]):
                p = PNG_DIR / f"{prop}.png"
                add_image_with_label(s, p, prop, Inches(1.0 + i * 4.0), Inches(5.7), side=Inches(0.9))
        except Exception:
            pass

    # 10. What didn't work
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Things We Tried That Didn't Work")
    add_bullets(s, [
        "Cross-modal hybrid (image kNN + curated descriptions, RRF-fused). jina-clip-v2's text-image alignment is weak on abstract glyphs. Rolled back.",
        "Multi-version cross-search (drop release_tag filter). More candidates → more chance of unrelated icons outranking the target. -9.6pp top-1.",
        "Mean-centering embeddings. -3.8pp; jina's outputs are already approximately centered.",
        "Query-time TTA. Neutral on top-1, +1.4pp top-10. Costs 4× embed calls.",
        "Each dead-end is a labeled commit so it's reproducible / restorable.",
    ])

    # 11. Tools & artifacts
    s = prs.slides.add_slide(BLANK)
    add_title(s, "What's Shipping")
    add_bullets(s, [
        "<IconSearch /> docs page component — paste, drop, or type",
        "MCP server — AI assistants can call icon_search over stdio",
        "Static HTML quality dashboard — sortable, filterable, color-coded view of every icon's rank",
        "Playwright quality sweep — repeatable evaluation against any vector field",
        "Background ingester — add new EUI versions in ~60 sec",
    ])

    # 12. Where the ceiling is
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Where the Ceiling Is Now")
    add_bullets(s, [
        "Current: 85.9% top-1 / 96.7% top-3 / 98.4% top-10 on v115",
        "Remaining failures are genuine semantic confusions the embedder can't see from visuals alone:",
        ("`stop` (hollow rounded square) ↔ `editorPositionBottomRight` (also a hollow rounded square)", 1),
        ("`beta` ↔ `editorBold`, `tokenString` ↔ `tokenSemanticText`", 1),
        "Three paths to higher accuracy, increasing cost:",
        ("LLM reranker over top-12 — projected +8-15pp top-1, ~$0.005/query, no training", 1),
        ("Trained cross-encoder on hard negatives — projected +5-10pp, 15-25h", 1),
        ("Click-to-confirm + log picks — every search becomes a labeled training pair", 1),
    ])

    # 13. Closing
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Wrap")
    add_bullets(s, [
        "Built a working multimodal icon-search tool grounded entirely in Elastic infrastructure.",
        "86% top-1 / 98% top-10 — without changing the model, just by treating \"what's in the index\" as a problem worth caring about.",
        "Repo: eui-embeddings + worktree feat/icon-vector-search on EUI",
        "Demo: paste an icon screenshot, see the result land.",
    ])

    prs.save(str(out_path))


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    parser.add_argument(
        "--out", type=Path,
        default=REPO_ROOT / "reports" / "eui_icon_search_deck.pptx",
    )
    args = parser.parse_args()
    logging.basicConfig(level=logging.INFO, format="%(message)s")
    args.out.parent.mkdir(parents=True, exist_ok=True)
    build_deck(args.out)
    log.info("wrote %s (%d KB)", args.out, args.out.stat().st_size // 1024)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
